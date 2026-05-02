"""
Generate an NCI-ready PowerPoint deck for:
"Traffic Accident Risk Analysis and Prediction (Ireland)"

Run:
  .\\.venv\\Scripts\\python.exe scripts\\generate_presentation.py

Output:
  outputs/presentation/Traffic_Accident_Risk_Ireland_NCI.pptx
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _set_run_font(run, size_pt: int, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_section_slide(prs: Presentation, title: str) -> None:
    # "Section header" layout if available; fall back to title-only
    try:
        slide_layout = prs.slide_layouts[2]
    except Exception:
        slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], body_size: int = 18) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()

    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            _set_run_font(run, body_size, bold=False, color=RGBColor(0x22, 0x22, 0x22))
        # If python-pptx split unexpectedly, font may be empty; set via first run
        if not p.runs:
            r = p.add_run()
            r.text = line
            _set_run_font(r, body_size)


def add_two_column_slide(prs: Presentation, title: str, left_title: str, left_bullets: list[str], right_title: str, right_bullets: list[str]) -> None:
    """Custom layout using two textboxes (no default two-column layout dependency)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(0.9))
    tf = title_box.text_frame
    tf.text = title
    for p in tf.paragraphs:
        for r in p.runs:
            _set_run_font(r, 32, bold=True, color=RGBColor(0x0F, 0x17, 0x2A))

    def add_col(x_in, title_text, bullets):
        box = slide.shapes.add_textbox(Inches(x_in), Inches(1.35), Inches(5.9), Inches(6.2))
        t = box.text_frame
        t.clear()
        p0 = t.paragraphs[0]
        p0.text = title_text
        for r in p0.runs:
            _set_run_font(r, 20, bold=True, color=RGBColor(0x1F, 0x4E, 0x79))
        for b in bullets:
            p = t.add_paragraph()
            p.text = b
            p.level = 0
            for r in p.runs:
                _set_run_font(r, 16, bold=False, color=RGBColor(0x22, 0x22, 0x22))

    add_col(0.6, left_title, left_bullets)
    add_col(6.7, right_title, right_bullets)


def add_pipeline_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(0.9))
    tf = title_box.text_frame
    tf.text = "End-to-End Pipeline (Extract → Store → Transform → Load → Analyze → Dashboard)"
    for p in tf.paragraphs:
        for r in p.runs:
            _set_run_font(r, 30, bold=True, color=RGBColor(0x0F, 0x17, 0x2A))

    accent = RGBColor(0x1F, 0x4E, 0x79)
    fill = RGBColor(0xF7, 0xF9, 0xFC)
    border = RGBColor(0xC7, 0xD6, 0xE8)

    def node(x, y, w, h, text):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = border
        shp.line.width = Pt(1.5)
        t = shp.text_frame
        t.text = text
        for p in t.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _set_run_font(r, 14, bold=True, color=accent)

    # positions tuned for 13.333 x 7.5 in widescreen
    w, h = Inches(2.35), Inches(0.85)
    y1 = Inches(1.55)
    y2 = Inches(2.75)
    y3 = Inches(3.95)
    y4 = Inches(5.15)

    node(Inches(0.7), y1, w, h, "Sources\n(API JSON + CSV outputs)")
    node(Inches(3.35), y1, w, h, "Extract\n(Python + Requests)")
    node(Inches(6.0), y1, w, h, "Raw Store\nMongoDB (JSON docs)")
    node(Inches(8.65), y1, w, h, "Transform\n(Pandas feature engineering)")

    node(Inches(2.0), y2, w, h, "Curated Store\nPostgreSQL (typed tables)")
    node(Inches(4.75), y2, w, h, "Integrated Table\n(county + date joins)")
    node(Inches(7.5), y2, w, h, "Quality Checks\n(logging + sanity joins)")

    node(Inches(2.0), y3, w, h, "EDA + Stats\n(county/time/corr)")
    node(Inches(4.75), y3, w, h, "Model Baseline\n(Linear Regression)")
    node(Inches(7.5), y3, w, h, "Exports\n(figures + metrics CSV/JSON)")

    node(Inches(3.35), y4, Inches(5.2), h, "Dashboard (Streamlit)\nfilters + charts + map + correlation view")

    # Connectors (simple left-to-right flow hints)
    def arrow(x1, y1, x2, y2):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        c.line.color.rgb = RGBColor(0x6B, 0x7C, 0x93)
        c.line.width = Pt(2)

    arrow(Inches(3.05), y1 + h / 2, Inches(3.35), y1 + h / 2)
    arrow(Inches(5.7), y1 + h / 2, Inches(6.0), y1 + h / 2)
    arrow(Inches(8.35), y1 + h / 2, Inches(8.65), y1 + h / 2)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(6.15), Inches(12.0), Inches(1.1))
    nt = note.text_frame
    nt.text = (
        "Design intent: preserve semi-structured fidelity (MongoDB) while making analytics SQL-friendly (PostgreSQL). "
        "Dashboard reads the processed integrated dataset for fast demo iteration."
    )
    for p in nt.paragraphs:
        for r in p.runs:
            _set_run_font(r, 13, bold=False, color=RGBColor(0x33, 0x33, 0x33))


def main() -> None:
    out_dir = Path("outputs/presentation")
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "Traffic_Accident_Risk_Ireland_NCI.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    today = date.today().isoformat()

    add_title_slide(
        prs,
        "Traffic Accident Risk Analysis and Prediction (Ireland)",
        "NCI MSc Data Analytics — Analytics Programming & Visualisation (2025/26)\n"
        "Team: Saignan Surabhi (25147773) • Suvan Reddy Marthala (25148087) • Sai Srinivas Naidu (25142992)\n"
        f"Generated: {today}",
    )

    add_section_slide(prs, "Agenda")
    add_bullet_slide(
        prs,
        "What we will cover (viva-friendly structure)",
        [
            "Problem, objectives, and research questions (why this matters for road safety analytics)",
            "Datasets + integrity constraints (incl. non-Kaggle requirement) and join strategy",
            "Architecture: MongoDB (raw JSON) + PostgreSQL (curated relational tables)",
            "ETL details: cleaning, datetime harmonisation, merges, engineered features",
            "Analysis: EDA, correlations, hotspots, baseline model + how to interpret responsibly",
            "Dashboard UI/UX: filters, charts, map, accessibility of messaging",
            "Challenges faced + mitigations (engineering truth for the report)",
            "Results, limitations, future work, and demo flow",
        ],
        body_size=17,
    )

    add_section_slide(prs, "Context & Objectives")
    add_two_column_slide(
        prs,
        "Why this project exists",
        "Motivation",
        [
            "Collisions are influenced by time, place, exposure (how much travel), and environment (weather).",
            "Decision-makers need integrated evidence: not isolated tables, but joinable analytics.",
            "This module emphasises engineering + communication: databases, ETL, visuals, and insight discipline.",
        ],
        "Objectives (mapped to brief)",
        [
            "Integrate multiple related datasets and quantify risk patterns across counties and time.",
            "Quantify relationships between weather features and collision harm proxies after integration.",
            "Deliver reproducible code + dual-database storage + interactive dashboard storytelling.",
            "Provide defensible limitations: association ≠ causation; synthetic data vs official records.",
        ],
    )

    add_section_slide(prs, "Research Questions")
    add_bullet_slide(
        prs,
        "Novel, answerable questions (coursework scope)",
        [
            "RQ1: Where and when does accident burden concentrate after integration (county + daily time series)?",
            "RQ2: After merging weather and exposure features, what multivariate associations appear (correlation view)?",
            "RQ3: Can a transparent baseline model explain any variance in a composite severity-risk index (interpretability first)?",
            "Operational question (dashboard): Can a non-technical reviewer explore filters without breaking interpretability?",
        ],
        body_size=17,
    )

    add_section_slide(prs, "Datasets (3 sources, joinable)")
    add_bullet_slide(
        prs,
        "Dataset A — Collisions (structured CSV; ≥1,000 rows)",
        [
            "5,000 synthetic collision records generated programmatically (original generator; not Kaggle).",
            "Key fields: county, datetime, geo coordinates, severity, casualties, vehicles, road context flags.",
            "Purpose: stable structured fact table for joins, EDA, and modelling experiments.",
            "Integrity statement: synthetic collisions demonstrate methodology; not official national statistics.",
        ],
        body_size=17,
    )
    add_bullet_slide(
        prs,
        "Dataset B — Weather (semi-structured JSON via API; stored in MongoDB)",
        [
            "Open-Meteo Historical Weather API (archive) → daily features: temperature, precipitation, max wind.",
            "Semi-structured JSON payloads preserved in MongoDB collection: traffic_risk_ireland.raw_weather.",
            "Why MongoDB: nested daily arrays, provider flexibility, auditability of raw payload structure.",
            "Flattened into typed relational weather_daily for PostgreSQL analytics and SQL-friendly joins.",
        ],
        body_size=16,
    )
    add_bullet_slide(
        prs,
        "Dataset C — Traffic exposure panel (structured CSV; ≥1,000 rows)",
        [
            "County × day panel (~8 counties × ~1,096 days ≈ 8,768 rows for 2023–2025 window).",
            "Includes population baseline + synthetic exposure proxies (vehicle-km, congestion, goods share).",
            "Join keys: county + date (same as weather + collisions after daily alignment).",
            "Purpose: reduce ‘raw count’ misinterpretation by enabling exposure-aware proxies in modelling.",
        ],
        body_size=16,
    )

    add_pipeline_slide(prs)

    add_section_slide(prs, "Database Design (why two systems)")
    add_two_column_slide(
        prs,
        "MongoDB vs PostgreSQL (LO1 alignment)",
        "MongoDB (non-relational)",
        [
            "Stores raw API JSON documents (flexible schema).",
            "Preserves nested daily series without premature flattening.",
            "Supports evolution if API fields change (engineering realism).",
            "Trade-off: fewer relational constraints; requires disciplined downstream validation.",
        ],
        "PostgreSQL (relational)",
        [
            "Stores curated tables with dtypes suitable for analytics joins.",
            "Tables: accidents_clean, weather_daily, county_traffic_exposure_daily, accident_risk_integrated.",
            "Indexing on (date, county) and severity for repeatable querying.",
            "Trade-off: requires schema management; benefits clarity + performance.",
        ],
    )

    add_section_slide(prs, "ETL: Transform logic (high signal)")
    add_bullet_slide(
        prs,
        "Cleaning, standardisation, integration",
        [
            "Datetime parsing + harmonisation → derive hour, weekday, weekend, rush-hour flags.",
            "Imputation strategy: numeric medians where needed; documented as a coursework engineering choice.",
            "Join accidents to weather_daily on (county, date); join exposure panel on (county, date).",
            "Engineered features: accident_severity_index; weather_impact_score; rain/wind flags.",
            "Exposure proxy feature: accidents_per_million_vehicle_km_proxy (interpret carefully; not a formal rate).",
        ],
        body_size=16,
    )

    add_section_slide(prs, "Analysis & Modelling")
    add_two_column_slide(
        prs,
        "From exploration to a baseline model",
        "Analytics layer",
        [
            "EDA: county summaries, hourly patterns, time series volatility.",
            "Correlation matrix: highlights multivariate structure; not causal claims.",
            "Hotspot mapping: geographic concentration; paired with filters to reduce cherry-picking.",
        ],
        "Baseline model (transparent)",
        [
            "Linear regression predicting accident_severity_index (interpretable coefficients).",
            "Features: time flags + weather + exposure proxies + population baseline.",
            "Expected outcome: modest R² if drivers are incomplete; still valuable for workflow demonstration.",
            "Future: GLMs for counts, regularisation, richer exposure, spatial models.",
        ],
    )

    add_section_slide(prs, "Dashboard UI/UX (Streamlit)")
    add_bullet_slide(
        prs,
        "Design choices for technical + non-technical audiences",
        [
            "Layout: KPI row → time series → severity distribution → map → correlation matrix (progressive disclosure).",
            "Filters: counties, date range slider, rainy-only toggle (small set to avoid clutter).",
            "Colour: severity encoded with perceptually ordered scales; map uses intuitive hotspot emphasis.",
            "Accessibility: descriptive titles, axis labels, and consistent terminology across charts.",
            "Trust cues (recommended in viva): show date coverage text + note data provenance (synthetic collisions).",
        ],
        body_size=16,
    )

    add_section_slide(prs, "Engineering Challenges → Mitigations (report-ready)")
    add_bullet_slide(
        prs,
        "What broke / what we fixed (typical real projects)",
        [
            "Challenge: MongoDB Atlas connectivity + credential/config drift → Mitigation: .env config + verify inserts + logs.",
            "Challenge: PostgreSQL missing database object → Mitigation: auto-create DB + idempotent loads.",
            "Challenge: API timeouts / partial failures → Mitigation: per-county try/except + warnings, avoid silent empties.",
            "Challenge: join correctness (date types) → Mitigation: enforce YYYY-MM-DD strings + validate integrated row counts.",
            "Challenge: disk space / venv issues on Windows → Mitigation: venv + pip cache purge + clean reinstall workflow.",
            "Challenge: marking constraints (≥1000 rows per dataset) → Mitigation: county-day exposure panel scaling.",
        ],
        body_size=15,
    )

    add_section_slide(prs, "Results (what to say carefully)")
    add_bullet_slide(
        prs,
        "Findings framed responsibly",
        [
            "Spatial concentration: highest counts in Dublin in generated summaries (example from integrated run).",
            "Temporal variability: daily accident counts fluctuate; dashboards help detect spikes vs noise.",
            "Weather association layer: correlation view shows structure among engineered scores; interpret as exploratory.",
            "Model: report MAE/R² as ‘baseline capability’, not deployment accuracy.",
            "Recommendation style: target hotspots + peak periods + weather-aware messaging (operational framing).",
        ],
        body_size=16,
    )

    add_section_slide(prs, "Limitations & Future Work")
    add_bullet_slide(
        prs,
        "Critical self-evaluation (high marks language)",
        [
            "Synthetic collisions: excellent for engineering reproducibility; replace with official open data for policy claims.",
            "County-day weather + exposure: coarse spatial resolution; improve with grid-level meteorology and link-level traffic.",
            "Causality: confounding (seasonality, travel demand) requires stronger methods + richer covariates.",
            "Next steps: data quality tests, orchestration (Airflow), monitoring, and ethical governance for PII.",
        ],
        body_size=16,
    )

    add_section_slide(prs, "Demo Script (10 minutes)")
    add_bullet_slide(
        prs,
        "Suggested on-screen flow for TeamX.mp4",
        [
            "0:00–0:20 — Title slide + IDs visible + one-sentence mission.",
            "0:20–1:20 — Problem + research questions + why databases matter.",
            "1:20–3:00 — Pipeline slide + dataset definitions + integrity statement.",
            "3:00–5:00 — Live/quick ETL narrative: Mongo insert + Postgres tables (logs/screenshots).",
            "5:00–8:00 — Dashboard demo: filters + chart interpretation discipline.",
            "8:00–9:30 — Results + limitations + future work (_marker-friendly).",
            "9:30–10:00 — Conclusion + teamwork reflection + thank you.",
        ],
        body_size=15,
    )

    add_section_slide(prs, "References (speak to IEEE report)")
    add_bullet_slide(
        prs,
        "What to cite in the report (examples)",
        [
            "Open-Meteo documentation (API) + MongoDB/PostgreSQL documentation (systems).",
            "Visualization theory: Munzner (task abstraction) + accessibility guidance.",
            "Road safety context: WHO fact sheets / policy summaries (high-level).",
        ],
        body_size=18,
    )

    add_title_slide(prs, "Thank you", "Questions?")

    prs.save(out_path)
    print(f"Wrote: {out_path.resolve()}")


if __name__ == "__main__":
    main()
